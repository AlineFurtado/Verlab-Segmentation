# -*- coding: utf-8 -*-
"""
Created on Tue Mar  5 20:59:29 2024

@author: User-Aline
"""

import cv2
from pptx import Presentation
from pptx.util import Inches, Pt
import numpy as np
import os
from io import BytesIO 
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN  # Importar PP_ALIGN para centralizar o texto

# Função para converter um array numpy em uma imagem codificada
def numpy_array_to_image(array):
    success, encoded_image = cv2.imencode('.jpg', array)
    if success:
        return encoded_image
    else:
        return None

# Função para carregar os frames de um vídeo
def load_video_frames(video_folder, extension, num_frames, target_size=(640, 480)):
    frames = []
    count = 0
    for filename in sorted(os.listdir(video_folder)):
        if filename.endswith('.' + extension):
            frame = cv2.imread(os.path.join(video_folder, filename))
            if frame is not None:
                frame = cv2.resize(frame, target_size)
                frames.append(frame)
                count += 1
                if count >= num_frames:  # Verificar se já carregamos o número desejado de frames
                    break
    print(f'{count} frames carregados')
    return frames

# Carregar os frames dos vídeos
video1_frames = load_video_frames('frames_cam1', 'jpg', 5400)


video4_frames = load_video_frames('frames_cam4', 'jpg', 5400)


# Carregar os frames dos plots
plots_frames = load_video_frames('plots_frames_cocktail', 'png', num_frames=120)


# Criar um arquivo PowerPoint
prs = Presentation()
prs.slide_width = Inches(13.333)  # Largura do slide em polegadas
prs.slide_height = Inches(7.5)    # Altura do slide em polegadas (correspondente ao widescreen 16:9)
print('arquivo ppt criado')

# Calcular o número de frames correspondentes a 3 segundos considerando a taxa de amostragem do vídeo
fps = 15  # Taxa de amostragem do vídeo (15 frames por segundo)
frames_por_segundo = 3  # Taxa de amostragem dos plots (um plot a cada 3 segundos)
frames_por_plot = fps * frames_por_segundo  # Número de frames correspondentes a 3 segundos

# Adicionar slides
for i in range(120):
    slide_layout = prs.slide_layouts[5]  # Seleciona o layout com dois conteúdos
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Social Space Segmentation for Approaching Tasks"

    
    # Alterar cor de fundo e texto
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0, 0, 0)  # Cor de fundo preto
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)  # Cor do texto branco
    
    # Adicionar caixa de texto para os autores
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))  # Ajuste de posição
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Aline F. F. Silva, Douglas G. Macharet"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(255, 255, 255)  # Cor do texto branco
    
    # Calcular o índice do frame correspondente ao plot atual
    index_video = i * (fps * frames_por_segundo)
    
    # Adicionar imagem do vídeo da camera1
    left_cam1 = Inches(0.193)  # Convertido para polegadas (0.49 cm)
    top_cam1 = Inches(2.989)    # Convertido para polegadas (7.58 cm)
    width = Inches(4)
    height = Inches(2.83)
    encoded_image1 = numpy_array_to_image(video1_frames[index_video])
    if encoded_image1 is not None:
        pic1 = slide.shapes.add_picture(BytesIO(encoded_image1), left_cam1, top_cam1, width, height)
    print(f'{video1_frames[index_video]} inserido no slide')
    
    # Adicionar legenda para o vídeo da camera1
    txBox1 = slide.shapes.add_textbox(left_cam1, Inches(5.7), width, Inches(0.5))  # Convertido para polegadas
    tf1 = txBox1.text_frame
    p1 = tf1.add_paragraph()
    p1.text = "Dataset Salsa CocktailParty camera1"
    p1.font.size = Pt(20)
    p1.font.color.rgb = RGBColor(255, 255, 255)  # Cor do texto branco
    txBox1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Centralizar texto
    
    # Adicionar imagem do vídeo da camera4
    left_cam4 = Inches(4.692)   # Convertido para polegadas (11.92 cm)
    encoded_image4 = numpy_array_to_image(video4_frames[index_video])
    if encoded_image4 is not None:
        pic2 = slide.shapes.add_picture(BytesIO(encoded_image4), left_cam4, top_cam1, width, height)
    print(f'{video4_frames[index_video]} inserido no slide')
    
    # Adicionar legenda para o vídeo da camera4
    txBox2 = slide.shapes.add_textbox(left_cam4, Inches(5.7), width, Inches(0.5))  # Convertido para polegadas
    tf2 = txBox2.text_frame
    p2 = tf2.add_paragraph()
    p2.text = "Dataset Salsa CocktailParty camera4"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(255, 255, 255)  # Cor do texto branco
    txBox2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Centralizar texto
    
    # Adicionar imagem do plot correspondente
    left_plot = Inches(9.157)   # Ajuste a posição horizontal conforme necessário
    encoded_plot = numpy_array_to_image(plots_frames[i])
    if encoded_plot is not None:
        pic3 = slide.shapes.add_picture(BytesIO(encoded_plot), left_plot, top_cam1, width, height)
    print(f'{plots_frames[i]} inserido no slide')
    
    # Adicionar legenda para o plot
    txBox3 = slide.shapes.add_textbox(left_plot, Inches(5.7), width, Inches(0.5))  # Convertido para polegadas
    tf3 = txBox3.text_frame
    p3 = tf3.add_paragraph()
    p3.text = "Methodology"  # Altere conforme necessário
    p3.font.size = Pt(20)
    p3.font.color.rgb = RGBColor(255, 255, 255)  # Cor do texto branco
    txBox3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Centralizar texto

    print(f'slide:{i} criado')

# Salvar o arquivo PowerPoint
prs.save('slides_test_1.pptx')
print('slides save')
